"""PPT generation service using python-pptx with template placeholder support."""
import os
from typing import Dict, Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

os.makedirs("output/ppts", exist_ok=True)
os.makedirs("output/previews", exist_ok=True)

_CONTENT_LEFT = Emu(457200)
_CONTENT_WIDTH = Emu(11277295)
_FOOTER_TEXT = "EduGuide AI Lesson System"


def _replace_in_runs(slide, replacements: Dict[str, str]) -> None:
    """Replace placeholder text in slide shapes while preserving font formatting."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for placeholder, value in replacements.items():
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, value)


def _delete_slide(prs: Presentation, index: int) -> None:
    rId = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


def _build_content_slide(slide, section: Dict[str, Any], page_num: int) -> None:
    """Build a content slide matching the template's layout positioning."""
    # Page number (top-left)
    num_box = slide.shapes.add_textbox(_CONTENT_LEFT, Emu(274320), Emu(1371600), Emu(365760))
    num_frame = num_box.text_frame
    num_frame.text = f"{page_num:02d}"
    num_frame.paragraphs[0].font.size = Pt(14)
    num_frame.paragraphs[0].font.bold = True

    # Section title
    title_box = slide.shapes.add_textbox(_CONTENT_LEFT, Emu(731520), _CONTENT_WIDTH, Emu(731520))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    section_title = section.get("title", "")
    duration = section.get("duration_minutes", 0)
    title_frame.text = f"{section_title}  ({duration}分钟)"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True

    # Content area
    content_box = slide.shapes.add_textbox(_CONTENT_LEFT, Emu(1645920), _CONTENT_WIDTH, Emu(4389120))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    activity = section.get("activity", "")
    method = section.get("teaching_method", "")
    outcome = section.get("expected_outcome", "")
    full_text = f"【教学活动】\n{activity}\n\n【教学方法】\n{method}\n\n【预期成果】\n{outcome}"
    content_frame.text = full_text
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(16)

    # Footer
    footer_box = slide.shapes.add_textbox(_CONTENT_LEFT, Emu(6400800), _CONTENT_WIDTH, Emu(274320))
    footer_frame = footer_box.text_frame
    footer_frame.text = _FOOTER_TEXT
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def _build_end_slide(slide, title: str = "教案") -> None:
    """Build an end slide matching the template style."""
    ty_box = slide.shapes.add_textbox(_CONTENT_LEFT, Emu(2286000), _CONTENT_WIDTH, Emu(1097280))
    ty_frame = ty_box.text_frame
    ty_frame.text = "感谢聆听！"
    ty_frame.paragraphs[0].font.size = Pt(40)
    ty_frame.paragraphs[0].font.bold = True
    ty_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    msg_box = slide.shapes.add_textbox(_CONTENT_LEFT, Emu(3657600), _CONTENT_WIDTH, Emu(548640))
    msg_frame = msg_box.text_frame
    msg_frame.text = f"{title} — Questions & Discussion"
    msg_frame.paragraphs[0].font.size = Pt(20)
    msg_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    line_box = slide.shapes.add_textbox(Emu(3657600), Emu(5943600), Emu(4876495), Emu(25400))
    line_box.text_frame.text = ""


def generate_ppt_from_lesson_plan(
    plan_data: Dict[str, Any],
    template: str = "lesson_default",
    ppt_id: str = "",
) -> str:
    """Generate PPT file from lesson plan data. Returns file_path."""
    template_path = f"templates/{template}.pptx"
    if os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    title = plan_data.get("title", "教案")
    objective = plan_data.get("objective", "")
    sections = plan_data.get("sections", [])
    has_template = len(prs.slides) >= 3

    if has_template:
        # Template has 3 slides: [0]=title, [1]=content_template, [2]=end
        # Step 1: Replace title slide placeholders (preserves template styling)
        _replace_in_runs(prs.slides[0], {
            "{{TITLE}}": title,
            "{{SUBTITLE}}": f"教学目标：{objective}",
        })

        # Step 2: Delete content template (slide 1) and end slide (slide 2)
        _delete_slide(prs, 2)
        _delete_slide(prs, 1)

        # Step 3: Add content slides for each section
        for i, section in enumerate(sections):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _build_content_slide(slide, section, i + 1)

        # Step 4: Add end slide
        end_slide = prs.slides.add_slide(prs.slide_layouts[6])
        _build_end_slide(end_slide, title)
    else:
        # No template: build everything from blank
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = title_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tb.text_frame.text = title
        tb.text_frame.paragraphs[0].font.size = Pt(44)
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        ob = title_slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
        ob.text_frame.text = f"教学目标：{objective}"
        ob.text_frame.paragraphs[0].font.size = Pt(20)
        ob.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for i, section in enumerate(sections):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _build_content_slide(slide, section, i + 1)

        end_slide = prs.slides.add_slide(prs.slide_layouts[6])
        _build_end_slide(end_slide, title)

    file_path = f"output/ppts/{ppt_id}.pptx"
    prs.save(file_path)
    return file_path


def get_slide_count(file_path: str) -> int:
    """Get number of slides in PPT."""
    prs = Presentation(file_path)
    return len(prs.slides)
